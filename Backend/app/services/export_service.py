import io
import tempfile
from typing import Tuple, Optional
from uuid import UUID
from datetime import datetime
import os

from fastapi import HTTPException, status
from fastapi.responses import StreamingResponse, FileResponse
from sqlalchemy.orm import Session
import docx
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE
from sqlalchemy import or_
from sqlalchemy import and_


from .. import models, schemas
from ..database import get_db

class ExportService:
    """Service for exporting documents and presentations."""
    
    def __init__(self, db: Session):
        self.db = db
    
    async def export_document(
        self,
        project_id: UUID,
        user_id: UUID,
        export_format: str = "docx",
        include_comments: bool = False,
        include_revision_history: bool = False
    ) -> Tuple[bytes, str, str]:
        """
        Export a project as a document.
        
        Args:
            project_id: The ID of the project to export
            user_id: The ID of the user requesting the export
            export_format: The format to export to (docx, pdf, txt)
            include_comments: Whether to include comments in the export
            include_revision_history: Whether to include revision history
            
        Returns:
            A tuple containing (file_content, content_type, filename)
        """
        # Verify project exists and belongs to user
        project = self.db.query(models.Project).filter(
            and_(
                models.Project.id == project_id,
                models.Project.user_id == user_id
            )
        ).first()
        
        if not project:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail="Project not found"
            )
        
        # Get all sections for the project
        sections = (
            self.db.query(models.Section)
            .filter(models.Section.project_id == project_id)
            .order_by(models.Section.idx)
            .all()
        )
        
        if not sections:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="No sections found in the project"
            )
        
        # Export based on format
        if export_format == "docx":
            return await self._export_to_docx(
                project=project,
                sections=sections,
                include_comments=include_comments,
                include_revision_history=include_revision_history
            )
        elif export_format == "pptx":
            return await self._export_to_pptx(
                project=project,
                sections=sections,
                include_comments=include_comments
            )
        elif export_format == "txt":
            return await self._export_to_txt(
                project=project,
                sections=sections,
                include_comments=include_comments
            )
        else:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Unsupported export format: {export_format}"
            )
    
    async def _export_to_docx(
        self,
        project: models.Project,
        sections: list[models.Section],
        include_comments: bool = False,
        include_revision_history: bool = False
    ) -> Tuple[bytes, str, str]:
        """Export project to a Word document."""
        # Create a new document
        doc = docx.Document()
        
        # Add title
        title = doc.add_heading(level=0)
        title_run = title.add_run(project.title)
        title_run.bold = True
        title_run.font.size = Pt(24)
        
        # Add creation date
        date_para = doc.add_paragraph()
        date_para.add_run(f"Created: {datetime.utcnow().strftime('%Y-%m-%d %H:%M:%S')} UTC")
        date_para.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        
        # Add sections
        for section in sections:
            # Add section title
            doc.add_heading(section.title, level=1)
            
            # Add section content
            if section.content:
                content_para = doc.add_paragraph()
                content_para.add_run(section.content)
            
            # Add comments if requested
            if include_comments and section.comments:
                doc.add_heading("Comments", level=2)
                for comment in section.comments:
                    comment_para = doc.add_paragraph(style='Intense Quote')
                    comment_para.add_run(f"{comment.user.email}: {comment.comment}")
            
            # Add revision history if requested
            if include_revision_history and section.revisions:
                doc.add_heading("Revision History", level=2)
                for revision in section.revisions:
                    rev_para = doc.add_paragraph()
                    rev_para.add_run(f"{revision.created_at.strftime('%Y-%m-%d %H:%M:%S')} - "
                                   f"{revision.user.email}\n")
                    rev_para.add_run(f"Prompt: {revision.prompt}\n").italic = True
                    rev_para.add_run(revision.generated_content)
                    doc.add_paragraph()  # Add space between revisions
            
            doc.add_page_break()
        
        # Save to a bytes buffer
        buffer = io.BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        
        filename = f"{project.title.replace(' ', '_')}.docx"
        return buffer.getvalue(), "application/vnd.openxmlformats-officedocument.wordprocessingml.document", filename
    
    async def _export_to_pptx(
        self,
        project: models.Project,
        sections: list[models.Section],
        include_comments: bool = False
    ) -> Tuple[bytes, str, str]:
        """Export project to a PowerPoint presentation."""
        # Create a new presentation
        prs = Presentation()
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = project.title
        subtitle.text = f"Generated on {datetime.utcnow().strftime('%Y-%m-%d')}"
        
        # Add section slides
        for section in sections:
            # Add section title slide
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = section.title
            
            # Add content
            if section.content:
                text_frame = content.text_frame
                text_frame.word_wrap = True
                text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                
                # Clear existing paragraphs (if any)
                text_frame.clear()
                
                # Split content into paragraphs
                for line in section.content.split("\n"):
                    if line.strip():
                        p = text_frame.add_paragraph()
                        p.text = line.strip()
                        p.font.size = Pt(18)
                        p.level = 0
            
            # Add comments if requested
            if include_comments and section.comments:
                # Add a new slide for comments
                slide_layout = prs.slide_layouts[1]
                comment_slide = prs.slides.add_slide(slide_layout)
                title = comment_slide.shapes.title
                content = comment_slide.placeholders[1]
                
                title.text = f"Comments: {section.title}"
                
                comments_text = "\n\n".join(
                    f"{comment.user.email}: {comment.comment}" 
                    for comment in section.comments
                )
                
                text_frame = content.text_frame
                text_frame.text = comments_text
        
        # Save to a bytes buffer
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        
        filename = f"{project.title.replace(' ', '_')}.pptx"
        return buffer.getvalue(), "application/vnd.openxmlformats-officedocument.presentationml.presentation", filename
    
    async def _export_to_txt(
        self,
        project: models.Project,
        sections: list[models.Section],
        include_comments: bool = False
    ) -> Tuple[bytes, str, str]:
        """Export project to a plain text file."""
        lines = []
        
        # Add title
        lines.append(f"{project.title}")
        lines.append("=" * len(project.title))
        lines.append(f"Generated on: {datetime.utcnow().strftime('%Y-%m-%d %H:%M:%S')} UTC")
        lines.append("")
        
        # Add sections
        for section in sections:
            lines.append(f"{section.title}")
            lines.append("-" * len(section.title))
            
            if section.content:
                lines.append(section.content)
            
            # Add comments if requested
            if include_comments and section.comments:
                lines.append("")
                lines.append("Comments:")
                lines.append("--------")
                for comment in section.comments:
                    lines.append(f"{comment.user.email}:")
                    lines.append(f"  {comment.comment}")
                    lines.append("")
            
            lines.append("\n")
        
        # Join lines and encode as UTF-8
        content = "\n".join(lines).encode('utf-8')
        filename = f"{project.title.replace(' ', '_')}.txt"
        
        return content, "text/plain; charset=utf-8", filename
